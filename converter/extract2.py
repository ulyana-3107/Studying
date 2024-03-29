from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import shutil
import aspose.slides as asp_slides
import aspose.pydrawing as drawing
import pptx
import json
from docx import Document
import argparse


def get_image_format(image_type):
    return {
        "jpeg": drawing.imaging.ImageFormat.jpeg,
        "emf": drawing.imaging.ImageFormat.emf,
        "bmp": drawing.imaging.ImageFormat.bmp,
        "png": drawing.imaging.ImageFormat.png,
        "wmf": drawing.imaging.ImageFormat.wmf,
        "gif": drawing.imaging.ImageFormat.gif,
    }.get(image_type, drawing.imaging.ImageFormat.jpeg)


def split_objects(pres_file: str) -> str:
    basename = str(Path().cwd()) + '\\' + 'data'
    if Path(basename).exists():
        shutil.rmtree(basename)
        Path(basename).mkdir()
    else:
        Path(basename).mkdir()

    pres1 = Presentation(pres_file)
    pres2 = asp_slides.Presentation(pres_file)

    n_text, n_image, n_links = 0, 0, 0
    sorted_objects, not_sorted_objects, sl_obj_nums = [], [], []
    im_dir, layouts_dir = basename + '\\' + 'images', basename + '\\' + 'layouts'
    links_dir = basename + '\\' + 'links'

    for i in range(len(pres1.slides)):
        sl, sl_elems = pres1.slides[i], []
        links, num = set(), 0

        for sh in sl.shapes:
            written = False
            if sh.has_text_frame:
                for par in sh.text_frame.paragraphs:
                    for run in par.runs:
                        address = run.hyperlink.address
                        if address is None:
                            continue

                        if not Path(links_dir).exists():
                            Path(links_dir).mkdir()
                        n_links += 1
                        num += 1
                        file_name = links_dir + '\\' + f'link_{n_links}'
                        with open(file_name, 'w', encoding='utf-8-sig') as file:
                            written = True
                            file.write(address)
                        links.add(str(address))
                        sl_elems.append([file_name, 'link'])

                if written:
                    continue

                dir_name = basename + '\\' + 'texts'
                if not Path(dir_name).exists():
                    Path(dir_name).mkdir()

                text_frame = sh.text_frame
                if text_frame.paragraphs:
                    num += 1
                    n_text += 1
                    f_name = f'{dir_name}\\text{n_text}.txt'
                    written = False
                    with open(f_name, 'w', encoding='utf-8-sig') as writer:
                        for paragraph in text_frame.paragraphs:
                                text = paragraph.text
                                if text and text not in links:
                                    writer.write(text + '\n')
                                    written = True
                                else:
                                    break
                    if not written:
                        Path(f_name).unlink()
                        n_text -= 1
                    else:
                        sl_elems.append([f_name, 'text'])

            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if not Path(im_dir).exists():
                    Path(im_dir).mkdir()
                n_image += 1
                image = sh.image
                image_bytes = image.blob
                f_name = im_dir + '\\' + f'image_{n_image}.{image.ext}'
                with open(f_name, 'wb') as file:
                    file.write(image_bytes)
                sl_elems.append([f_name, 'image'])

            # extracting shapes
            # elif sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # extracting group shapes
            # elif sh.shape_type == MSO_SHAPE_TYPE.GROUP:

        sorted_objects.append(sl_elems)
        sl_obj_nums.append(len(sl.shapes) - num)

    # Extracting of images
    if len(pres2.images):
        if not Path(im_dir).exists():
            Path(im_dir).mkdir()
        for image in pres2.images:
            n_image += 1
            file_name = "image_{0}.{1}"
            im_type = image.content_type.split("/")[1]
            im_format = get_image_format(im_type)
            image.system_image.save(im_dir + '\\' + file_name.format(n_image, im_type), im_format)
            file = im_dir + '\\' + file_name.format(n_image, im_type) + f'.{im_format}'
            not_sorted_objects.append([file, 'image'])

    # extracting of background pictures
    slide_index = 0
    for slide in pres2.slides:
        slide_index += 1
        image_format = drawing.imaging.ImageFormat.jpeg
        back_image = None
        file_name = "BackImage_Slide_{0}{1}.{2}"
        is_layout = False

        if slide.background.fill_format.fill_type == asp_slides.FillType.PICTURE:
            back_image = slide.background.fill_format.picture_fill_format.picture.image
        elif slide.layout_slide.background.fill_format.fill_type == asp_slides.FillType.PICTURE:
            back_image = slide.layout_slide.background.fill_format.picture_fill_format.picture.image
            is_layout = True

        if back_image is not None:
            if not Path(layouts_dir).exists():
                Path(layouts_dir).mkdir()
            image_type = back_image.content_type.split("/")[1]
            im_format = get_image_format(image_type)
            file = layouts_dir + '\\' + file_name.format("layout_" if is_layout else "", slide_index, image_type) + f'.{im_format}'
            back_image.system_image.save(file)
            not_sorted_objects.append([file, 'slide_layout'])

    return write_json_data(sorted_objects)


def write_json_data(sorted_objects: list) -> str:
    data = {}

    for i in range(len(sorted_objects)):
        name = f'slide_{i + 1}'
        slide_data, slide = {}, sorted_objects[i]
        for elem in slide:
            slide_data[elem[0]] = elem[1]
        data[name] = slide_data

    file_name = 'pres_objects.json'
    with open(file_name, 'w') as objects:
        json.dump(data, objects, indent=3)

    return file_name


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('pres', type=str, help='Path to presentation')
    args = parser.parse_args()
    split_objects(args.pres)



